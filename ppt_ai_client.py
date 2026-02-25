import os
import time
import requests
import logging

logger = logging.getLogger(__name__)

# ==============================
# Configuración
# ==============================

PPTAI_API_KEY = os.getenv("PPTAI_API_KEY")
USE_MOCK = os.getenv("USE_MOCK_APIS", "false").lower() == "true"

BASE_URL = "https://api.ppt.ai/v1"

HEADERS = {
    "Authorization": f"Bearer {PPTAI_API_KEY}",
    "Content-Type": "application/json"
}

# ==============================
# API pública
# ==============================

def crear_presentacion_desde_texto(
    texto: str,
    archivo_salida: str,
    timeout: int = 300
) -> str:
    """
    Crea una presentación PPTX a partir de texto.
    Devuelve la ruta al archivo generado.
    """

    # -------- MOCK --------
    if USE_MOCK:
        logger.info("[MOCK] ppt.ai: creando presentación simulada")
        _crear_ppt_falso(archivo_salida)
        return archivo_salida

    # -------- REAL --------
    if not PPTAI_API_KEY:
        raise RuntimeError("PPTAI_API_KEY no configurada")

    payload = {
        "title": "Resumen automático",
        "language": "es",
        "content": texto,
        "slides": {
            "min": 5,
            "max": 15
        },
        "style": "education"
    }

    logger.info("Creando presentación con ppt.ai...")
    r = requests.post(
        f"{BASE_URL}/presentations",
        headers=HEADERS,
        json=payload,
        timeout=30
    )
    r.raise_for_status()

    job_id = r.json()["id"]
    logger.info("Job ppt.ai creado: %s", job_id)

    return _esperar_y_descargar(job_id, archivo_salida, timeout)

# ==============================
# Funciones internas
# ==============================

def _esperar_y_descargar(job_id: str, archivo_salida: str, timeout: int) -> str:
    inicio = time.time()

    while True:
        r = requests.get(
            f"{BASE_URL}/presentations/{job_id}",
            headers=HEADERS,
            timeout=30
        )
        r.raise_for_status()

        data = r.json()
        estado = data.get("status")

        if estado == "completed":
            download_url = data["download_url"]
            return _descargar(download_url, archivo_salida)

        if estado == "failed":
            raise RuntimeError(f"ppt.ai falló: {data}")

        if time.time() - inicio > timeout:
            raise TimeoutError("Timeout esperando ppt.ai")

        time.sleep(5)

def _descargar(url: str, destino: str) -> str:
    logger.info("Descargando PPTX desde ppt.ai...")
    r = requests.get(url, stream=True, timeout=60)
    r.raise_for_status()

    with open(destino, "wb") as f:
        for chunk in r.iter_content(8192):
            f.write(chunk)

    logger.info("Presentación guardada en %s", destino)
    return destino

def _crear_ppt_falso(destino: str):
    """
    Genera un PPTX mínimo para modo MOCK
    """
    from pptx import Presentation

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Presentación simulada"
    slide.placeholders[1].text = "ppt.ai MOCK"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Contenido"
    slide.placeholders[1].text = "Esta presentación es un mock."

    prs.save(destino)
